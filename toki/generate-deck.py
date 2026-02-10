#!/usr/bin/env python3
"""
TokiStorage Partnership Deck Generator
Generates PPTX files (JP + EN) and converts to PDF via LibreOffice.
Design aligned with TokiStorage landing page (index.css).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import subprocess, os, sys, re, zipfile, io

# ── Design tokens (matched to index.css :root) ────────────────────────
TOKI_BLUE     = RGBColor(0x25, 0x63, 0xEB)  # --toki-blue
TOKI_BLUE_DK  = RGBColor(0x1D, 0x4E, 0xD8)  # --toki-blue-dark
TOKI_BLUE_PALE= RGBColor(0xEF, 0xF6, 0xFF)  # --toki-blue-pale
GOLD          = RGBColor(0xC9, 0xA9, 0x62)  # --toki-gold
GOLD_PALE     = RGBColor(0xFD, 0xF8, 0xE8)  # --toki-gold-pale
EMERALD       = RGBColor(0x16, 0xA3, 0x4A)  # green-600 (check marks in LP)
GREEN_PALE    = RGBColor(0xF0, 0xFD, 0xF4)  # green-50

TEXT_PRIMARY  = RGBColor(0x1E, 0x29, 0x3B)  # --text-primary
TEXT_SECONDARY= RGBColor(0x47, 0x55, 0x69)  # --text-secondary
TEXT_MUTED    = RGBColor(0x94, 0xA3, 0xB8)  # --text-muted
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)

BG_PAGE       = RGBColor(0xF8, 0xFA, 0xFC)  # --bg-page
BG_SECTION    = RGBColor(0xF1, 0xF5, 0xF9)  # --bg-section
BORDER        = RGBColor(0xE2, 0xE8, 0xF0)  # --border

# Dark backgrounds (LP hero / footer)
DARK_BG       = RGBColor(0x1E, 0x29, 0x3B)  # --text-primary as bg (LP footer)
DARK_BG2      = RGBColor(0x0F, 0x17, 0x2A)  # LP footer gradient start

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)  # 16:9

FONT_JP = "IPAPGothic"
FONT_EN = "Calibri"

OUT_DIR = os.path.dirname(os.path.abspath(__file__))


# ── Helpers ────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_rect(slide, left, top, width, height, fill=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def set_text(tf, text, font_name, size, color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_para(tf, text, font_name, size, color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_textbox(slide, left, top, width, height, text, font_name, size,
                color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        txBox.text_frame.paragraphs[0].space_before = Pt(0)
        txBox.text_frame.paragraphs[0].space_after = Pt(0)
    except:
        pass
    from pptx.oxml.ns import qn
    txBody = txBox.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if anchor == MSO_ANCHOR.MIDDLE:
        bodyPr.set('anchor', 'ctr')
    elif anchor == MSO_ANCHOR.BOTTOM:
        bodyPr.set('anchor', 'b')
    set_text(tf, text, font_name, size, color, bold, align)
    return txBox


# ── Action bar ─────────────────────────────────────────────────────────

def add_action_bar(slide, text, font):
    bar_h = Inches(0.55)
    add_rect(slide, 0, 0, SLIDE_W, bar_h, fill=DARK_BG)
    add_textbox(slide, Inches(0.5), 0, SLIDE_W - Inches(1), bar_h,
                text, font, 9, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)


# ── Footer ─────────────────────────────────────────────────────────────

def add_footer(slide, left_text, pg, font):
    y = SLIDE_H - Inches(0.35)
    add_rect(slide, 0, y, SLIDE_W, Pt(0.5), fill=BORDER)
    add_textbox(slide, Inches(0.5), y + Pt(2), Inches(4), Inches(0.28),
                left_text, font, 6.5, TEXT_MUTED)
    add_textbox(slide, Inches(4), y + Pt(2), Inches(2), Inches(0.28),
                "Confidential", font, 6.5, TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, SLIDE_W - Inches(1), y + Pt(2), Inches(0.5), Inches(0.28),
                str(pg), font, 6.5, TEXT_MUTED, bold=True, align=PP_ALIGN.RIGHT)


# ── Section label ──────────────────────────────────────────────────────

def add_section_label(slide, text, font, top):
    add_textbox(slide, Inches(0.5), top, Inches(3), Inches(0.25),
                text.upper(), font, 7, TOKI_BLUE, bold=True)


# ── Card helpers (all use single border style: BORDER, 0.75pt) ────────

def draw_col_card(slide, x, y, w, h, num, title, body, font):
    add_rect(slide, x, y, w, h, fill=WHITE, border_color=BORDER)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.08), Inches(0.5), Inches(0.25),
                num, font, 11, TOKI_BLUE, bold=True)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.35), w - Inches(0.24), Inches(0.22),
                title, font, 8, TEXT_PRIMARY, bold=True)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.58), w - Inches(0.24), h - Inches(0.68),
                body, font, 7, TEXT_SECONDARY)


def draw_grid_card(slide, x, y, w, h, icon_letter, title, body, font):
    add_rect(slide, x, y, w, h, fill=WHITE, border_color=BORDER)
    # Icon circle
    ix, iy = x + Inches(0.12), y + Inches(0.12)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, ix, iy, Inches(0.32), Inches(0.32))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TOKI_BLUE_PALE
    circ.line.fill.background()
    add_textbox(slide, ix, iy, Inches(0.32), Inches(0.32),
                icon_letter, font, 9, TOKI_BLUE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, x + Inches(0.52), y + Inches(0.12), w - Inches(0.65), Inches(0.2),
                title, font, 8, TEXT_PRIMARY, bold=True)
    add_textbox(slide, x + Inches(0.52), y + Inches(0.34), w - Inches(0.65), h - Inches(0.42),
                body, font, 6.5, TEXT_SECONDARY)


def draw_model_item(slide, x, y, w, h, badge_text, badge_color, title, body, example, font):
    add_rect(slide, x, y, w, h, fill=WHITE, border_color=BORDER)
    # Badge
    bx, by, bw, bh = x + Inches(0.12), y + Inches(0.1), Inches(0.55), Inches(0.55)
    badge = add_rect(slide, bx, by, bw, bh, fill=badge_color)
    add_textbox(slide, bx, by, bw, bh, badge_text, font, 6, WHITE, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx = x + Inches(0.8)
    tw = w - Inches(0.95)
    add_textbox(slide, tx, y + Inches(0.08), tw, Inches(0.2),
                title, font, 7.5, TEXT_PRIMARY, bold=True)
    add_textbox(slide, tx, y + Inches(0.28), tw, Inches(0.25),
                body, font, 6.5, TEXT_SECONDARY)
    add_textbox(slide, tx, y + Inches(0.52), tw, Inches(0.2),
                example, font, 6, TOKI_BLUE)


def draw_flow_box(slide, x, y, w, h, title, body, bg_color, title_color, font):
    add_rect(slide, x, y, w, h, fill=bg_color, border_color=BORDER)
    add_textbox(slide, x + Inches(0.08), y + Inches(0.06), w - Inches(0.16), Inches(0.22),
                title, font, 8, title_color, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.08), y + Inches(0.3), w - Inches(0.16), h - Inches(0.36),
                body, font, 6.5, TEXT_SECONDARY, align=PP_ALIGN.CENTER)


def draw_sector_card(slide, x, y, w, h, icon_letter, title, body, font):
    add_rect(slide, x, y, w, h, fill=WHITE, border_color=BORDER)
    # Icon circle
    circ_x = x + (w - Inches(0.36)) / 2
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, circ_x, y + Inches(0.08), Inches(0.36), Inches(0.36))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TOKI_BLUE_PALE
    circ.line.fill.background()
    add_textbox(slide, circ_x, y + Inches(0.08), Inches(0.36), Inches(0.36),
                icon_letter, font, 9, TOKI_BLUE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, x + Inches(0.06), y + Inches(0.48), w - Inches(0.12), Inches(0.2),
                title, font, 7, TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.06), y + Inches(0.68), w - Inches(0.12), h - Inches(0.73),
                body, font, 6, TEXT_SECONDARY, align=PP_ALIGN.CENTER)


def draw_table_row(slide, y, row_h, cells, font, is_header=False, is_first_col_header=False):
    col_widths = [Inches(1.3), Inches(3.6), Inches(3.7)]
    x = Inches(0.7)
    for i, (text, width) in enumerate(zip(cells, col_widths)):
        bg = DARK_BG if is_header else (BG_SECTION if (i == 0 and is_first_col_header) else WHITE)
        fg = WHITE if is_header else (TEXT_PRIMARY if (i == 0 and is_first_col_header) else TEXT_SECONDARY)
        bld = is_header or (i == 0 and is_first_col_header)
        fs = 6.5 if is_header else 7
        rect = add_rect(slide, x, y, width, row_h, fill=bg, border_color=BORDER, border_width=Pt(0.5))
        if i == 2 and not is_header:
            rect.fill.solid()
            rect.fill.fore_color.rgb = TOKI_BLUE_PALE
            fg = TOKI_BLUE_DK
        add_textbox(slide, x + Inches(0.08), y, width - Inches(0.16), row_h,
                    text, font, fs, fg, bold=bld, anchor=MSO_ANCHOR.MIDDLE)
        x += width


# ══════════════════════════════════════════════════════════════════════
#  CONTENT DATA
# ══════════════════════════════════════════════════════════════════════

CONTENT = {
    "ja": {
        "font": FONT_JP,
        "filename": "tokistorage-partnership-deck",
        "cover": {
            "label": "Partnership Proposal \u2014 Confidential",
            "title": "貴社クライアントに、\n「千年の存在証明」を。",
            "sub": "テクノロジーは揃っています。ユースケースも200以上。\n足りないのは、届ける仕組みです。",
            "org": "Universal Need株式会社",
            "product": "TokiStorage",
        },
        "s2": {
            "bar": "既存のデジタルサービスは「今」に最適化されており、千年スケールの保存レイヤーが構造的に不在である",
            "label": "Background",
            "cards": [
                ("01", "デジタルは「今」に最適化",
                 "クラウドやSNSは日常の記録に極めて優秀。しかし100年・1000年スケールは設計対象外。時間軸が違えば、必要な設計も違う。"),
                ("02", "「墓じまい」が社会課題に",
                 "無縁墓は年間数万基。2040年には単身世帯が4割超。「誰が記憶を残すのか」は個人の問題から社会の問題へ移行。"),
                ("03", "千年レイヤーが不在",
                 "デジタル遺品整理、AI故人再現——市場は急成長。しかし全て既存インフラ上。補完する千年レイヤーが求められている。"),
            ],
            "footer": "TokiStorage \u2014 協業提案",
        },
        "s3": {
            "bar": "TokiStorageは技術・ユースケース・思想基盤・展開網の4つを「設計思想ごと」提供できる",
            "label": "Our Offering",
            "cards": [
                ("Q", "石英ガラス記録技術",
                 "金属蒸着によるQRコード刻印。サーバー・電源ゼロ。SLA 100%、1000年保証。スマホカメラで読取可能。"),
                ("U", "200+ ユースケース（業界別整理済み）",
                 "終活・婚礼・寺社・学校・企業・自治体・NGO・ホテル・航空。提案書にそのまま転用可能な粒度。"),
                ("E", "70+ 思想エッセイ（9領域）",
                 "存在証明を心理学・宗教・経済・AI・宇宙まで展開。知的コンテンツとしてクライアント提案に活用可能。"),
                ("A", "Pearl Soap + アンバサダー網",
                 "贈与経済の実践。全国にワークショップ展開可能な分散型運動体。エンドユーザーとの接点を提供。"),
            ],
            "footer": "TokiStorage \u2014 協業提案",
        },
        "s4": {
            "bar": "TokiStorageは既存デジタルサービスの「競合」ではなく「千年レイヤー」として補完する位置づけである",
            "label": "Positioning",
            "headers": ["", "デジタルサービス（日常の記録）", "TokiStorage（千年の記録）"],
            "rows": [
                ("記録媒体", "クラウド / HDD", "石英ガラス（物理）"),
                ("得意な時間軸", "今〜数十年（日常利用に最適）", "100年〜1000年（永続保存に特化）"),
                ("インフラ", "サーバー・電源（利便性の源泉）", "不要（GitHub分散管理）"),
                ("読み取り", "専用アプリ / ログイン", "スマホのカメラだけ"),
                ("文化的厚み", "機能・利便性が価値の中心", "70+ エッセイ＋贈与経済の実践"),
                ("社会的接点", "プラットフォームとしての貢献", "SoulCarrier（無縁墓・遺骨送還活動）"),
            ],
            "footer": "TokiStorage \u2014 協業提案",
        },
        "s5": {
            "bar": "貴社のビジネスモデルに応じた3つの提携モデルを用意しており、段階的な移行も可能",
            "label": "Partnership Models",
            "models": [
                ("紹介\nモデル", TOKI_BLUE,
                 "A. クライアント紹介型パートナーシップ",
                 "紹介ベースで連携。リファラルフィーをお支払い。クライアント対応・納品はTokiStorage側で完結。",
                 "例：終活コンサル→存在証明を提案 / 葬祭業DX→メモリアルオプション追加"),
                ("共同\n提案", GOLD,
                 "B. 共同ソリューション型",
                 "貴社コンサルにTokiStorageを組み込んだ共同提案。ESG・地方創生・文化保存の「出口」として千年記録を位置づけ。",
                 "例：自治体DX→地域記憶アーカイブ / ホテルCX改革→ゲスト記録のアップグレード"),
                ("事業\n共創", EMERALD,
                 "C. 新規事業共創型",
                 "存在証明を軸に新規事業を共同立ち上げ。技術・思想・ユースケースはTokiStorage、市場アクセス・信用・スケールは貴社。",
                 "例：メモリアルテック新規事業 / 企業向け永続アーカイブサービス"),
            ],
            "footer": "TokiStorage \u2014 協業提案",
        },
        "s6": {
            "bar": "クライアント・パートナー・TokiStorageの三者がWinになる収益設計を採用している",
            "label": "Revenue Flow",
            "flows": [
                ("クライアント", "千年の存在証明\n社会的意義の実感", GREEN_PALE, EMERALD),
                ("パートナー（貴社）", "紹介フィー or 共同提案収益\nクライアントLTV向上", TOKI_BLUE_PALE, TOKI_BLUE),
                ("TokiStorage", "技術・思想・納品\n収益の一部→SoulCarrier活動", GOLD_PALE, RGBColor(0x92, 0x40, 0x0E)),
            ],
            "callout_title": "初期パートナー優遇",
            "callout_body": "複数のコンサルティングファームに順次ご提案を進めています。最初に提携いただいたファームには、紹介モデルの優先条件・エリア独占権など、初期パートナーならではのインセンティブをご用意します。",
            "footer": "TokiStorage \u2014 協業提案",
        },
        "s7": {
            "bar": "葬祭・ホスピタリティ・宗教法人・自治体・ESG・金融の6領域で特に高い親和性がある",
            "label": "Client Fit",
            "lead": "貴社のクライアントポートフォリオに、以下のセクターはありませんか。",
            "sectors": [
                ("M", "葬祭・メモリアル", "墓じまい代替、永代供養デジタル化、遺族向け新サービス"),
                ("H", "ホスピタリティ", "ウェディング記録、ホテルCX、記念日サービス"),
                ("R", "宗教法人・寺社", "檀家記録の永続化、参拝体験DX、文化財保存"),
                ("G", "自治体・教育", "地域アーカイブ、災害記録、学校史の永続化"),
                ("E", "ESG・サステナビリティ", "1000年設計の企業理念記録、SDGs実績の永続証明"),
                ("F", "金融・保険", "終活関連サービス連携、デジタル遺品対策"),
            ],
            "footer": "TokiStorage \u2014 協業提案",
        },
        "s8": {
            "bar": "代表はBig4出身であり、ファームのコンプライアンス要件を理解した上で提携モデルを設計している",
            "label": "Team & Independence",
            "name": "佐藤卓也 \u2014 Universal Need株式会社 代表",
            "bio": "大手コンサルティングファームでの経験を経て、半導体製造装置のエンジニアリング20年超。タイムレスタウン新浦安（250世帯）の自治会長として「ゆりかごから墓場まで」のコミュニティ運営を経験。SoulCarrier活動で「記憶が消える恐怖」を目の当たりにし、TokiStorageを着想。マウイ・山中湖でのオフグリッド実証を経て、制度に依存しない千年設計の技術を完成。",
            "tags": ["元Big4ファーム", "半導体エンジニアリング 20年+", "自治会長（250世帯）",
                     "SoulCarrier主宰", "オフグリッド実証済み", "佐渡島移住予定（2026春）"],
            "ind_title": "独立性（Independence）について",
            "ind_body": "本提携はベンダーパートナーシップです。SalesforceやSAPの導入推奨と同じ構造であり、監査契約・出資関係は一切含みません。独立性に関する懸念が発生しない設計です。",
            "footer": "TokiStorage \u2014 協業提案",
        },
        "s9": {
            "bar": "Next Step",
            "title": "まずは30分、お話ししませんか。",
            "sub": "提携の形は柔軟に設計できます。\n貴社クライアントの具体的な課題から逆算して、\n一緒に最適なモデルを見つけましょう。",
            "footer_left": "Universal Need株式会社 \u2014 TokiStorage",
        },
    },
    "en": {
        "font": FONT_EN,
        "filename": "tokistorage-partnership-deck-en",
        "cover": {
            "label": "Partnership Proposal \u2014 Confidential",
            "title": "Bring your clients\n1,000-year proof of existence.",
            "sub": "The technology is ready. Over 200 use cases mapped.\nWhat's missing is the delivery network.",
            "org": "Universal Need Inc.",
            "product": "TokiStorage",
        },
        "s2": {
            "bar": "Existing digital services are optimized for \"now\" \u2014 the millennium preservation layer is structurally absent",
            "label": "Background",
            "cards": [
                ("01", "Digital is optimized for \"now\"",
                 "Cloud and social platforms excel at everyday recording. But 100- or 1,000-year preservation is outside their design scope. Different time horizons need different architectures."),
                ("02", "End-of-life is now societal",
                 "Tens of thousands of graves go unclaimed yearly. By 2040, single-person households will exceed 40%. \"Who preserves memory?\" is now a social question."),
                ("03", "The millennium layer is missing",
                 "Digital estate management, AI recreations \u2014 the market is booming. Yet every solution runs on existing infrastructure. A complementary millennium layer is needed."),
            ],
            "footer": "TokiStorage \u2014 Partnership Proposal",
        },
        "s3": {
            "bar": "TokiStorage delivers technology, use cases, intellectual foundation, and distribution as a unified design philosophy",
            "label": "Our Offering",
            "cards": [
                ("Q", "Quartz glass recording",
                 "QR codes inscribed via metal deposition. Zero servers, zero power. SLA 100%, guaranteed 1,000 years. Readable by any smartphone camera."),
                ("U", "200+ use cases (organized by industry)",
                 "End-of-life, weddings, temples, schools, corporations, municipalities, NGOs, hotels, airlines. Ready for direct proposal integration."),
                ("E", "70+ philosophical essays (9 domains)",
                 "Proof of existence explored across psychology, religion, economics, AI, and space. Standalone intellectual content for client proposals."),
                ("A", "Pearl Soap + Ambassador network",
                 "A gift-economy practice and decentralized workshop network ready to scale nationwide. Direct end-user touchpoint."),
            ],
            "footer": "TokiStorage \u2014 Partnership Proposal",
        },
        "s4": {
            "bar": "TokiStorage is not a \"competitor\" to digital services \u2014 it is a complementary millennium layer",
            "label": "Positioning",
            "headers": ["", "Digital services (everyday records)", "TokiStorage (millennium records)"],
            "rows": [
                ("Medium", "Cloud / HDD", "Quartz glass (physical)"),
                ("Best horizon", "Now to decades (optimized for daily use)", "100\u20131,000 years (optimized for permanence)"),
                ("Infrastructure", "Servers & power (source of convenience)", "None required (GitHub distributed)"),
                ("Reading", "App / login required", "Any smartphone camera"),
                ("Cultural depth", "Functionality & convenience at the core", "70+ essays + gift economy practice"),
                ("Social impact", "Platform-level contribution", "SoulCarrier (unclaimed graves mission)"),
            ],
            "footer": "TokiStorage \u2014 Partnership Proposal",
        },
        "s5": {
            "bar": "Three partnership models tailored to your business model, with progressive escalation possible",
            "label": "Partnership Models",
            "models": [
                ("Referral", TOKI_BLUE,
                 "A. Client Referral Partnership",
                 "Introduce clients when TokiStorage fits. You receive a referral fee; we handle delivery end-to-end.",
                 "E.g.: End-of-life consulting \u2192 offer proof of existence / Funeral DX \u2192 add memorial option"),
                ("Joint", GOLD,
                 "B. Joint Solution Partnership",
                 "Embed TokiStorage into your consulting engagements. Position millennium records as the \"outcome layer\" of ESG, revitalization, or DX projects.",
                 "E.g.: Municipal DX \u2192 community archive / Hotel CX \u2192 guest record upgrade"),
                ("Co-Create", EMERALD,
                 "C. New Business Co-Creation",
                 "Launch a new venture together. We bring technology, philosophy, and use cases. You bring market access, credibility, and scale.",
                 "E.g.: Memorial-tech startup / Enterprise perpetual archive service"),
            ],
            "footer": "TokiStorage \u2014 Partnership Proposal",
        },
        "s6": {
            "bar": "Revenue design ensures all three parties \u2014 client, partner, and TokiStorage \u2014 win",
            "label": "Revenue Flow",
            "flows": [
                ("Client", "1,000-year proof of existence\nTangible social meaning", GREEN_PALE, EMERALD),
                ("Partner (you)", "Referral fee or joint revenue\nClient LTV increase", TOKI_BLUE_PALE, TOKI_BLUE),
                ("TokiStorage", "Technology & delivery\nRevenue \u2192 SoulCarrier mission", GOLD_PALE, RGBColor(0x92, 0x40, 0x0E)),
            ],
            "callout_title": "Early Partner Advantage",
            "callout_body": "We are approaching consulting firms sequentially. The first firm to partner receives preferential terms \u2014 including priority referral conditions and potential regional exclusivity. Early movers shape the partnership.",
            "footer": "TokiStorage \u2014 Partnership Proposal",
        },
        "s7": {
            "bar": "Six client sectors show particularly high affinity: funeral, hospitality, religious, government, ESG, and finance",
            "label": "Client Fit",
            "lead": "Does your client portfolio include any of these sectors?",
            "sectors": [
                ("M", "Funeral & Memorial", "Gravestone alternatives, digital perpetual care, bereavement services"),
                ("H", "Hospitality", "Wedding records, hotel CX, anniversary services"),
                ("R", "Religious Institutions", "Perpetual congregation records, visitor DX, cultural preservation"),
                ("G", "Government & Education", "Community archives, disaster records, school history"),
                ("E", "ESG & Sustainability", "1,000-year corporate purpose records, SDG impact proof"),
                ("F", "Finance & Insurance", "End-of-life service integration, digital estate"),
            ],
            "footer": "TokiStorage \u2014 Partnership Proposal",
        },
        "s8": {
            "bar": "The founder is a Big Four alumnus who designed the partnership model with full awareness of firm compliance",
            "label": "Team & Independence",
            "name": "Takuya Sato \u2014 CEO, Universal Need Inc.",
            "bio": "Former Big Four consultant \u2014 understands firm culture, client engagement, and project design from the inside. 20+ years in semiconductor manufacturing engineering. Former president of Timeless Town Shin-Urayasu residents' association (250 households). Through SoulCarrier's work with unclaimed graves, witnessed firsthand how memories vanish \u2014 and conceived TokiStorage. Validated off-grid, institution-free 1,000-year design through testing in Maui and Lake Yamanakako.",
            "tags": ["Big Four Alumni", "Semiconductor engineering 20+ yrs", "Community president (250 households)",
                     "SoulCarrier founder", "Off-grid validated", "Relocating to Sado Island (Spring 2026)"],
            "ind_title": "A note on independence",
            "ind_body": "This is a vendor partnership \u2014 structurally identical to recommending Salesforce or SAP. No audit engagement, no equity relationship, no independence concerns. Designed with full awareness of firm compliance requirements.",
            "footer": "TokiStorage \u2014 Partnership Proposal",
        },
        "s9": {
            "bar": "Next Step",
            "title": "Let's start with a\n30-minute conversation.",
            "sub": "Partnership structure is flexible by design.\nLet's reverse-engineer the right model\nfrom your clients' actual challenges.",
            "footer_left": "Universal Need Inc. \u2014 TokiStorage",
        },
    },
}


# ══════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════

def build_cover(prs, d):
    slide = add_blank_slide(prs)
    font = d["font"]
    c = d["cover"]
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=DARK_BG)
    add_textbox(slide, Inches(1), Inches(0.8), Inches(5), Inches(0.3),
                c["label"], font, 7, TEXT_MUTED, align=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.2),
                c["title"], font, 22, WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(2.9), Inches(7), Inches(0.8),
                c["sub"], font, 10, RGBColor(0xBB, 0xBB, 0xCC))
    # Bottom accent line
    stripe_y = SLIDE_H - Inches(0.6)
    add_rect(slide, 0, stripe_y, SLIDE_W, Inches(0.03), fill=TOKI_BLUE)
    add_textbox(slide, Inches(0.5), stripe_y + Inches(0.08), Inches(3), Inches(0.3),
                c["org"], font, 6.5, TEXT_MUTED)
    add_textbox(slide, Inches(4), stripe_y + Inches(0.08), Inches(2), Inches(0.3),
                c["product"], font, 6.5, TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, SLIDE_W - Inches(2), stripe_y + Inches(0.08), Inches(1.5), Inches(0.3),
                "Confidential", font, 6.5, TEXT_MUTED, align=PP_ALIGN.RIGHT)


def build_slide2(prs, d):
    """Problem: 3 column cards"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s2"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.7))
    card_w = Inches(2.85)
    card_h = Inches(3.6)
    gap = Inches(0.23)
    start_x = Inches(0.5)
    y = Inches(1.0)
    for i, (num, title, body) in enumerate(s["cards"]):
        x = start_x + i * (card_w + gap)
        draw_col_card(slide, x, y, card_w, card_h, num, title, body, font)
    add_footer(slide, s["footer"], 2, font)


def build_slide3(prs, d):
    """Solution: 2x2 grid"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s3"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.7))
    card_w = Inches(4.3)
    card_h = Inches(1.7)
    gap_x = Inches(0.25)
    gap_y = Inches(0.2)
    start_x = Inches(0.5)
    start_y = Inches(1.0)
    for i, (icon, title, body) in enumerate(s["cards"]):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        draw_grid_card(slide, x, y, card_w, card_h, icon, title, body, font)
    add_footer(slide, s["footer"], 3, font)


def build_slide4(prs, d):
    """Differentiator: comparison table"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s4"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.7))
    y = Inches(1.0)
    row_h = Inches(0.38)
    draw_table_row(slide, y, row_h, s["headers"], font, is_header=True)
    y += row_h
    for cells in s["rows"]:
        draw_table_row(slide, y, row_h, cells, font, is_first_col_header=True)
        y += row_h
    add_footer(slide, s["footer"], 4, font)


def build_slide5(prs, d):
    """Partnership models: 3 rows"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s5"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.7))
    item_w = Inches(8.6)
    item_h = Inches(0.75)
    gap = Inches(0.12)
    start_y = Inches(1.0)
    x = Inches(0.5)
    for i, (badge, color, title, body, ex) in enumerate(s["models"]):
        y = start_y + i * (item_h + gap)
        draw_model_item(slide, x, y, item_w, item_h, badge, color, title, body, ex, font)
    add_footer(slide, s["footer"], 5, font)


def build_slide6(prs, d):
    """Revenue flow + callout"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s6"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.7))
    box_w = Inches(2.4)
    box_h = Inches(0.85)
    gap = Inches(0.5)
    total = 3 * box_w + 2 * gap
    start_x = (SLIDE_W - total) / 2
    y = Inches(1.2)
    for i, (title, body, bg, tc) in enumerate(s["flows"]):
        x = start_x + i * (box_w + gap)
        draw_flow_box(slide, x, y, box_w, box_h, title, body, bg, tc, font)
        if i < 2:
            arrow_x = x + box_w + Inches(0.1)
            add_textbox(slide, arrow_x, y + Inches(0.2), Inches(0.3), Inches(0.3),
                        "\u2190", font, 14, TEXT_MUTED, align=PP_ALIGN.CENTER)
    # Callout (left accent stripe only, no outer border)
    cx, cy = Inches(0.5), Inches(2.4)
    cw, ch = Inches(8.6), Inches(1.0)
    add_rect(slide, cx, cy, cw, ch, fill=TOKI_BLUE_PALE)
    add_rect(slide, cx, cy, Inches(0.05), ch, fill=TOKI_BLUE)
    add_textbox(slide, cx + Inches(0.2), cy + Inches(0.08), cw - Inches(0.3), Inches(0.22),
                s["callout_title"], font, 7.5, TEXT_PRIMARY, bold=True)
    add_textbox(slide, cx + Inches(0.2), cy + Inches(0.32), cw - Inches(0.3), ch - Inches(0.4),
                s["callout_body"], font, 6.5, TEXT_SECONDARY)
    add_footer(slide, s["footer"], 6, font)


def build_slide7(prs, d):
    """Client sectors: 3x2 grid"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s7"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.7))
    add_textbox(slide, Inches(0.5), Inches(0.95), Inches(8), Inches(0.25),
                s["lead"], font, 8, TEXT_SECONDARY)
    card_w = Inches(2.85)
    card_h = Inches(1.45)
    gap_x = Inches(0.23)
    gap_y = Inches(0.18)
    start_x = Inches(0.5)
    start_y = Inches(1.3)
    for i, (icon, title, body) in enumerate(s["sectors"]):
        col = i % 3
        row = i // 3
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        draw_sector_card(slide, x, y, card_w, card_h, icon, title, body, font)
    add_footer(slide, s["footer"], 7, font)


def build_slide8(prs, d):
    """Founder + independence"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s8"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.7))
    # Avatar circle
    ax, ay = Inches(0.5), Inches(1.0)
    avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, ax, ay, Inches(0.6), Inches(0.6))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = DARK_BG
    avatar.line.fill.background()
    initials = "佐" if d is CONTENT["ja"] else "TS"
    add_textbox(slide, ax, ay, Inches(0.6), Inches(0.6),
                initials, font, 12, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(1.25), Inches(1.0), Inches(7), Inches(0.25),
                s["name"], font, 8, TEXT_PRIMARY, bold=True)
    add_textbox(slide, Inches(1.25), Inches(1.28), Inches(7.5), Inches(0.95),
                s["bio"], font, 6.5, TEXT_SECONDARY)
    # Tags
    tag_x = Inches(1.25)
    tag_y = Inches(2.2)
    for tag in s["tags"]:
        tw = Inches(len(tag) * 0.085 + 0.25)
        add_rect(slide, tag_x, tag_y, tw, Inches(0.22), fill=BG_SECTION, border_color=BORDER)
        add_textbox(slide, tag_x + Inches(0.05), tag_y, tw - Inches(0.1), Inches(0.22),
                    tag, font, 5.5, TEXT_SECONDARY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tag_x += tw + Inches(0.08)
        if tag_x > Inches(8.5):
            tag_x = Inches(1.25)
            tag_y += Inches(0.28)
    # Independence callout (left accent stripe only, no outer border)
    cy = Inches(2.6)
    cw, ch = Inches(8.6), Inches(0.8)
    cx = Inches(0.5)
    add_rect(slide, cx, cy, cw, ch, fill=TOKI_BLUE_PALE)
    add_rect(slide, cx, cy, Inches(0.05), ch, fill=TOKI_BLUE)
    add_textbox(slide, cx + Inches(0.2), cy + Inches(0.06), cw - Inches(0.3), Inches(0.2),
                s["ind_title"], font, 7, TEXT_PRIMARY, bold=True)
    add_textbox(slide, cx + Inches(0.2), cy + Inches(0.28), cw - Inches(0.3), ch - Inches(0.35),
                s["ind_body"], font, 6.5, TEXT_SECONDARY)
    add_footer(slide, s["footer"], 8, font)


def build_slide9(prs, d):
    """CTA slide"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s9"]
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=DARK_BG)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(3), Inches(0.35),
                s["bar"], font, 9, WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.0),
                s["title"], font, 18, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(2.8), Inches(7), Inches(1.0),
                s["sub"], font, 9, RGBColor(0xBB, 0xBB, 0xCC), align=PP_ALIGN.CENTER)
    # Footer
    stripe_y = SLIDE_H - Inches(0.5)
    add_rect(slide, 0, stripe_y - Inches(0.03), SLIDE_W, Pt(0.5), fill=RGBColor(0x33, 0x44, 0x55))
    add_textbox(slide, Inches(0.5), stripe_y, Inches(4), Inches(0.3),
                s["footer_left"], font, 6.5, TEXT_MUTED)
    add_textbox(slide, Inches(4), stripe_y, Inches(2), Inches(0.3),
                "Confidential", font, 6.5, TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, SLIDE_W - Inches(1), stripe_y, Inches(0.5), Inches(0.3),
                "9", font, 6.5, TEXT_MUTED, bold=True, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════

def generate(lang):
    d = CONTENT[lang]
    prs = new_prs()
    build_cover(prs, d)
    build_slide2(prs, d)
    build_slide3(prs, d)
    build_slide4(prs, d)
    build_slide5(prs, d)
    build_slide6(prs, d)
    build_slide7(prs, d)
    build_slide8(prs, d)
    build_slide9(prs, d)

    pptx_path = os.path.join(OUT_DIR, f"{d['filename']}.pptx")
    prs.save(pptx_path)
    _strip_theme_shadows(pptx_path)
    print(f"  PPTX saved: {pptx_path}")
    return pptx_path


def _strip_theme_shadows(pptx_path):
    """Remove default outerShdw / 3D effects from theme1.xml."""
    clean = '<a:effectStyleLst>' + \
            '<a:effectStyle><a:effectLst/></a:effectStyle>' * 3 + \
            '</a:effectStyleLst>'
    buf = io.BytesIO()
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.namelist():
                data = zin.read(item)
                if item == 'ppt/theme/theme1.xml':
                    text = data.decode('utf-8')
                    text = re.sub(r'<a:effectStyleLst>.*?</a:effectStyleLst>',
                                  clean, text, flags=re.DOTALL)
                    data = text.encode('utf-8')
                zout.writestr(item, data)
    with open(pptx_path, 'wb') as f:
        f.write(buf.getvalue())


def convert_to_pdf(pptx_path):
    out_dir = os.path.dirname(pptx_path)
    env = os.environ.copy()
    env["HOME"] = "/tmp"
    result = subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", out_dir, pptx_path],
        capture_output=True, text=True, timeout=120, env=env
    )
    if result.returncode != 0:
        print(f"  ERROR: {result.stderr}")
        return None
    pdf_path = pptx_path.replace(".pptx", ".pdf")
    print(f"  PDF saved: {pdf_path}")
    return pdf_path


if __name__ == "__main__":
    print("=== TokiStorage Partnership Deck Generator ===\n")

    for lang in ["ja", "en"]:
        print(f"[{lang.upper()}] Generating PPTX...")
        pptx = generate(lang)
        print(f"[{lang.upper()}] Converting to PDF...")
        convert_to_pdf(pptx)

    print("\nDone!")
